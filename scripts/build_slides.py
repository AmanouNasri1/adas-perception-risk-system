"""
Build the V7 presentation deck (docs/slides.pptx) — 8 slides following the project-plan
structure (PDF section 15.4). Reproducible: regenerate any time with

    python scripts/build_slides.py

Figures are pulled from outputs/figures/. Numbers mirror the README Results section.
"""
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

FIG = Path("outputs/figures")

# ADAS / automotive palette — dark navy sandwich with an amber "warning" accent.
NAVY  = RGBColor(0x0E, 0x1B, 0x2C)
PANEL = RGBColor(0x16, 0x27, 0x3D)
ICE   = RGBColor(0xCA, 0xDC, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
AMBER = RGBColor(0xFF, 0xB0, 0x00)
TEAL  = RGBColor(0x4F, 0xC3, 0xD9)
INK   = RGBColor(0x14, 0x20, 0x30)
MUTE  = RGBColor(0x6B, 0x7B, 0x8C)
LIGHT = RGBColor(0xF4, 0xF7, 0xFA)

HEAD_FONT = "Trebuchet MS"
BODY_FONT = "Calibri"

EMU_IN = 914400
SW, SH = 13.333, 7.5


def _solid(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def bg(slide, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(SW), Inches(SH))
    _solid(s, color)
    slide.shapes._spTree.remove(s._element)
    slide.shapes._spTree.insert(2, s._element)
    return s


def text(slide, l, t, w, h, runs, size=16, color=INK, bold=False, font=BODY_FONT,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line=1.05, space_after=4):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(2)
    items = runs if isinstance(runs, list) else [(runs, {})]
    for i, (txt, ov) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = ov.get("align", align)
        p.line_spacing = ov.get("line", line)
        p.space_after = Pt(ov.get("space_after", space_after))
        p.space_before = Pt(ov.get("space_before", 0))
        r = p.add_run()
        r.text = txt
        r.font.name = ov.get("font", font)
        r.font.size = Pt(ov.get("size", size))
        r.font.bold = ov.get("bold", bold)
        r.font.color.rgb = ov.get("color", color)
    return tb


def bullets(slide, l, t, w, h, items, size=16, color=INK, gap=7, marker_color=AMBER):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        p.line_spacing = 1.04
        head, _, rest = it.partition("|")
        r = p.add_run(); r.text = "▸ "
        r.font.name = BODY_FONT; r.font.size = Pt(size); r.font.bold = True
        r.font.color.rgb = marker_color
        r1 = p.add_run(); r1.text = head.strip()
        r1.font.name = BODY_FONT; r1.font.size = Pt(size); r1.font.bold = bool(rest)
        r1.font.color.rgb = color
        if rest:
            r2 = p.add_run(); r2.text = "  " + rest.strip()
            r2.font.name = BODY_FONT; r2.font.size = Pt(size); r2.font.color.rgb = color
    return tb


def title(slide, txt, color=INK, top=0.55, size=34):
    text(slide, 0.7, top, SW - 1.4, 1.0, txt, size=size, color=color, bold=True, font=HEAD_FONT)


def image_fit(slide, path, l, t, maxw, maxh, align="center", valign="middle"):
    iw, ih = Image.open(path).size
    ar = iw / ih
    w, h = maxw, maxw / ar
    if h > maxh:
        h, w = maxh, maxh * ar
    if align == "center":
        l = l + (maxw - w) / 2
    elif align == "right":
        l = l + (maxw - w)
    if valign == "middle":
        t = t + (maxh - h) / 2
    elif valign == "bottom":
        t = t + (maxh - h)
    return slide.shapes.add_picture(str(path), Inches(l), Inches(t), Inches(w), Inches(h))


def stat(slide, l, t, w, number, label, num_color=AMBER, lab_color=INK, num_size=40):
    text(slide, l, t, w, 0.8, number, size=num_size, color=num_color, bold=True,
         font=HEAD_FONT, align=PP_ALIGN.CENTER)
    text(slide, l, t + 0.72, w, 0.6, label, size=12, color=lab_color,
         align=PP_ALIGN.CENTER, line=1.0)


def rrect(slide, l, t, w, h, fill, radius=0.08):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t),
                               Inches(w), Inches(h))
    _solid(s, fill)
    try:
        s.adjustments[0] = radius
    except Exception:
        pass
    s.shadow.inherit = False
    return s


def caption(slide, l, t, w, txt, color=MUTE):
    text(slide, l, t, w, 0.3, txt, size=10.5, color=color, align=PP_ALIGN.CENTER, line=1.0)


def build():
    prs = Presentation()
    prs.slide_width = Inches(SW)
    prs.slide_height = Inches(SH)
    blank = prs.slide_layouts[6]

    # ---- Slide 1: Title / Problem & motivation (dark, hero image right) ----
    s = prs.slides.add_slide(blank); bg(s, NAVY)
    rrect(s, 7.05, 0, 6.3, SH, PANEL, radius=0.0)
    image_fit(s, FIG / "v6_lane_demo.png", 7.25, 1.7, 5.9, 4.1)
    caption(s, 7.25, 5.85, 5.9, "Live pipeline: drivable-area zone, lane lines, tracked vehicles, warnings", ICE)
    text(s, 0.7, 0.7, 6.1, 0.5, "BACHELOR-THESIS FOUNDATION  ·  COMPUTER VISION", size=13, color=AMBER, bold=True, font=HEAD_FONT)
    text(s, 0.7, 1.5, 6.2, 2.4, "Real-Time ADAS Perception & Risk Warning", size=40, color=WHITE, bold=True, font=HEAD_FONT, line=1.02)
    text(s, 0.7, 3.7, 6.1, 1.4,
         "A monocular, camera-only pipeline that turns dashcam video into actionable "
         "collision warnings — built as clean, measured engineering.", size=16, color=ICE, line=1.15)
    for i, (n, lab) in enumerate([("6", "phases shipped\nV1–V6"),
                                  ("0.922", "KITTI\nmAP@0.5"),
                                  ("34", "FPS\nbaseline")]):
        stat(s, 0.7 + i * 2.0, 5.35, 1.9, n, lab, num_color=AMBER, lab_color=ICE, num_size=34)
    text(s, 0.7, 6.95, 6.0, 0.4, "Amanou Allah Nasri", size=12, color=MUTE)

    # ---- Slide 2: System architecture ----
    s = prs.slides.add_slide(blank); bg(s, WHITE)
    title(s, "System architecture", color=INK)
    text(s, 0.7, 1.35, 11.9, 0.5,
         "One pipeline, not disconnected scripts — reusable package + thin CLI entrypoints, config-driven.",
         size=15, color=MUTE)
    stages = [("Frame", "cv2 loader"), ("Detect", "YOLO11s\n(KITTI)"), ("Track", "ByteTrack"),
              ("Lane / road", "YOLOPv2 seg"), ("Risk zones", "drivable +\nped/cyclist"),
              ("Distance", "f·H / h"), ("TTC", "closing speed"), ("Warnings", "tiered engine")]
    n = len(stages); gapx = 0.22
    bw = (SW - 1.4 - gapx * (n - 1)) / n
    by, bh = 2.5, 1.5
    for i, (h1, h2) in enumerate(stages):
        x = 0.7 + i * (bw + gapx)
        box = rrect(s, x, by, bw, bh, NAVY if i % 2 == 0 else PANEL, radius=0.12)
        tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = h1; r.font.bold = True; r.font.size = Pt(13.5)
        r.font.color.rgb = WHITE; r.font.name = HEAD_FONT
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run(); r2.text = h2; r2.font.size = Pt(10); r2.font.color.rgb = ICE
        r2.font.name = BODY_FONT
        if i < n - 1:
            text(s, x + bw - 0.02, by + bh / 2 - 0.25, gapx + 0.04, 0.5, "→",
                 size=18, color=AMBER, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, 0.7, 4.4, 11.9, 0.4, "Outputs", size=15, color=AMBER, bold=True, font=HEAD_FONT)
    bullets(s, 0.7, 4.9, 6.0, 2.2, [
        "Annotated demo video | color-coded boxes, zones, HUD",
        "warnings.csv | frame, track, class, warning, distance_m, ttc_s, bbox",
        "tracking_results.csv | persistent IDs per frame",
    ], size=14)
    bullets(s, 6.9, 4.9, 5.7, 2.2, [
        "frame_timings.csv | per-stage latency incl. lane_ms",
        "Figures & metrics | FPS plot, histograms, comparisons",
        "Config-driven | thresholds & models in configs/*.yaml",
    ], size=14)

    # ---- Slide 3: Detection and tracking ----
    s = prs.slides.add_slide(blank); bg(s, WHITE)
    title(s, "Detection & tracking pipeline", color=INK)
    bullets(s, 0.7, 1.5, 5.5, 3.6, [
        "Detector | YOLO11s, fine-tuned on KITTI (50 epochs, Colab T4)",
        "Accuracy | mAP@0.5 = 0.922, mAP@0.5:0.95 = 0.700",
        "Precision / Recall | 0.924 / 0.855",
        "Tracking | ByteTrack default; persistent track IDs",
        "Compared | BoT-SORT — fewer IDs (115 vs 126), longer tracks",
        "Classes | car, truck, bus, motorcycle (+ person/cyclist)",
    ], size=15, gap=9)
    for i, (n, lab) in enumerate([("0.922", "mAP@0.5"), ("0.70", "mAP@.5:.95"), ("116m", "train time")]):
        stat(s, 0.75 + i * 1.85, 5.25, 1.8, n, lab)
    rrect(s, 6.5, 1.5, 6.1, 4.6, LIGHT, radius=0.04)
    image_fit(s, FIG / "kitti_training_curve.png", 6.7, 1.7, 5.7, 2.0)
    image_fit(s, FIG / "tracker_comparison.png", 6.7, 3.85, 5.7, 2.05)
    caption(s, 6.5, 6.15, 6.1, "KITTI fine-tuning curve (top) · ByteTrack vs BoT-SORT (bottom)")

    # ---- Slide 4: ADAS warning logic ----
    s = prs.slides.add_slide(blank); bg(s, WHITE)
    title(s, "ADAS warning logic", color=INK)
    text(s, 0.7, 1.45, 5.7, 0.5, "Tiered warnings by severity (Priority 5 = highest):",
         size=15, color=MUTE)
    tiers = [("5", "TTC WARNING", "imminent collision, TTC < 3 s", AMBER),
             ("4", "PEDESTRIAN / CYCLIST", "vulnerable user in risk zone", RGBColor(0xE0,0x3C,0x3C)),
             ("3", "VEHICLE TOO CLOSE", "large bbox, near bottom of frame", RGBColor(0xF0,0x80,0x20)),
             ("1", "FRONT OBJECT", "object in ego drivable lane", TEAL)]
    y = 2.1
    for pr, name, desc, col in tiers:
        c = slide_circle = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), Inches(y), Inches(0.55), Inches(0.55))
        _solid(c, col)
        ctf = c.text_frame; ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
        cp = ctf.paragraphs[0]; cp.alignment = PP_ALIGN.CENTER
        cr = cp.add_run(); cr.text = pr; cr.font.bold = True; cr.font.size = Pt(18)
        cr.font.color.rgb = WHITE; cr.font.name = HEAD_FONT
        text(s, 1.45, y - 0.04, 4.9, 0.4, name, size=15, color=INK, bold=True, font=HEAD_FONT)
        text(s, 1.45, y + 0.32, 4.9, 0.4, desc, size=12, color=MUTE)
        y += 0.78
    bullets(s, 0.7, 5.5, 5.7, 1.6, [
        "Distance | D = f_y · H_real / h_bbox  (heuristic or KITTI-calibrated)",
        "TTC | distance / closing-speed, least-squares over a history buffer",
    ], size=13.5)
    rrect(s, 6.6, 1.5, 6.0, 4.9, LIGHT, radius=0.04)
    image_fit(s, FIG / "v5_ttc_frame.png", 6.8, 1.7, 5.6, 4.4)
    caption(s, 6.6, 6.5, 6.0, "HUD with distance, per-track TTC, and active warnings")

    # ---- Slide 5: Results and demo screenshots ----
    s = prs.slides.add_slide(blank); bg(s, WHITE)
    title(s, "Results & demo", color=INK)
    for i, (n, lab, c) in enumerate([("34", "FPS baseline\n(RTX 3050)", AMBER),
                                     ("15.6", "FPS w/ lane\nsegmentation", AMBER),
                                     ("98%", "lane-detection\nrate", TEAL),
                                     ("0.922", "detector\nmAP@0.5", TEAL)]):
        stat(s, 0.75 + i * 3.0, 1.45, 2.8, n, lab, num_color=c, num_size=40)
    rrect(s, 0.7, 3.0, 6.05, 4.0, LIGHT, radius=0.04)
    image_fit(s, FIG / "v6_lane_demo.png", 0.9, 3.2, 5.65, 3.4)
    caption(s, 0.7, 6.6, 6.05, "Full pipeline — V6 lane-aware front zone")
    rrect(s, 6.95, 3.0, 5.65, 4.0, LIGHT, radius=0.04)
    image_fit(s, FIG / "fps_plot.png", 7.15, 3.5, 5.25, 2.9)
    caption(s, 6.95, 6.6, 5.65, "Per-frame processing FPS over the clip")

    # ---- Slide 6: Evaluation metrics and failure cases ----
    s = prs.slides.add_slide(blank); bg(s, WHITE)
    title(s, "Evaluation & failure cases", color=INK)
    text(s, 0.7, 1.4, 11.9, 0.5,
         "Every design change is justified by measured A/B warning statistics:", size=15, color=MUTE)
    bullets(s, 0.7, 2.05, 5.7, 2.4, [
        "FRONT OBJECT false positives | −67%  (3637 → 1215) via lane-aware zone",
        "TTC false positives | −75%  (6181 → 1559) via in-path gating",
        "Other tiers unchanged | ped/cyclist & too-close logic isolated",
        "Honest cost | +1 deep model/frame: 34 → 15.6 FPS",
    ], size=14, gap=10)
    text(s, 0.7, 4.6, 5.7, 0.4, "Documented failure mode", size=15, color=AMBER, bold=True, font=HEAD_FONT)
    text(s, 0.7, 5.05, 5.7, 1.6,
         "At intersections there is no clear ego lane, so the drivable-area zone is "
         "low-confidence and the system falls back to the static front trapezoid "
         "(shown right) — a deliberate, logged behaviour.", size=13.5, color=INK, line=1.12)
    rrect(s, 6.7, 1.5, 5.9, 2.55, LIGHT, radius=0.04)
    image_fit(s, FIG / "v6_static_vs_lane_4200.png", 6.85, 1.65, 2.6, 2.25, align="center")
    image_fit(s, FIG / "v6_fallback_intersection.png", 9.5, 1.95, 2.95, 1.65)
    caption(s, 6.7, 4.1, 5.9, "Static vs lane zone (left) · intersection fallback (right)")
    rrect(s, 6.7, 4.55, 5.9, 2.1, LIGHT, radius=0.04)
    image_fit(s, FIG / "confidence_hist.png", 6.9, 4.7, 5.5, 1.8)
    caption(s, 6.7, 6.55, 5.9, "Detection-confidence distribution")

    # ---- Slide 7: Roadmap to thesis extension ----
    s = prs.slides.add_slide(blank); bg(s, WHITE)
    title(s, "Roadmap to thesis extension", color=INK)
    done = [("V1", "Detection + tracking + ADAS demo + metrics"),
            ("V2", "Evaluation: histograms, tracker comparison, failures"),
            ("V3", "KITTI fine-tuning (mAP@0.5 = 0.922)"),
            ("V4", "Monocular distance estimation"),
            ("V5", "Time-to-collision warnings"),
            ("V6", "Lane/road-area awareness (YOLOPv2)")]
    y = 1.55
    for tag, desc in done:
        c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), Inches(y), Inches(0.42), Inches(0.42))
        _solid(c, TEAL)
        ct = c.text_frame; ct.vertical_anchor = MSO_ANCHOR.MIDDLE
        cp = ct.paragraphs[0]; cp.alignment = PP_ALIGN.CENTER
        cr = cp.add_run(); cr.text = "✓"; cr.font.bold = True; cr.font.size = Pt(15)
        cr.font.color.rgb = WHITE
        text(s, 1.3, y - 0.02, 1.0, 0.4, tag, size=14, color=INK, bold=True, font=HEAD_FONT)
        text(s, 2.2, y - 0.02, 4.3, 0.45, desc, size=12.5, color=MUTE)
        y += 0.62
    rrect(s, 6.9, 1.5, 5.7, 5.1, NAVY, radius=0.05)
    text(s, 7.2, 1.8, 5.1, 0.5, "Now → thesis", size=18, color=AMBER, bold=True, font=HEAD_FONT)
    bullets(s, 7.2, 2.45, 5.2, 4.0, [
        "V7 (now) | technical report, slides, CV bullets, thesis proposal",
        "Depth study | heuristic vs calibrated vs learned, error by range",
        "Risk zones | learned vs fixed, in-path precision/recall vs labels",
        "TTC validation | vs ground-truth closing speed on KITTI tracking",
        "Deployment | fp16 / ONNX, accuracy-vs-FPS operating point",
    ], size=13.5, color=ICE, marker_color=AMBER, gap=11)

    # ---- Slide 8: Conclusion (dark) ----
    s = prs.slides.add_slide(blank); bg(s, NAVY)
    text(s, 0.9, 1.5, 11.5, 0.5, "CONCLUSION", size=15, color=AMBER, bold=True, font=HEAD_FONT)
    text(s, 0.9, 2.15, 11.5, 2.0,
         "A complete, reproducible monocular ADAS perception pipeline — and a measured "
         "foundation for a bachelor thesis.", size=30, color=WHITE, bold=True, font=HEAD_FONT, line=1.05)
    bullets(s, 0.9, 4.2, 11.5, 1.8, [
        "Six shipped phases | detection → tracking → lane seg → distance → TTC → warnings",
        "Measured, not claimed | every FPS and false-positive reduction is A/B-quantified",
        "Thesis-ready | clean repo, configs, CSV logs, documented failure modes",
    ], size=15, color=ICE, marker_color=AMBER, gap=10)
    text(s, 0.9, 6.6, 11.5, 0.4,
         "github.com/AmanouNasri1/adas-perception-risk-system", size=13, color=TEAL, bold=True)

    out = Path("docs/slides.pptx")
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    print(f"wrote {out}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")


if __name__ == "__main__":
    build()
